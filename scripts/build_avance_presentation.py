#!/usr/bin/env python3
"""Build Benchmark avance examples.pptx from the taxonomy template + draft items.

Uses raw episode images (what a VLM would see), not annotated GT frames.
"""

from __future__ import annotations

import argparse
import json
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt
from PIL import Image

# --- paths ---
REPO = Path(__file__).resolve().parents[1]
TEMPLATE = Path("/home/andreina/Documents/Programs/Benchmark - avance.pptx")
OUTPUT = Path("/home/andreina/Documents/Programs/Benchmark - avance examples.pptx")
DEFAULT_DRAFT_JSONS = [
    REPO / "src/cm_benchmark/storage/ai2thor/items/draft_house_007514.json",
    REPO / "src/cm_benchmark/storage/ai2thor/items/draft_house_001030.json",
]

# --- palette (from template theme) ---
NAVY = RGBColor(0x00, 0x2F, 0x4A)
TEAL = RGBColor(0x00, 0x93, 0x84)
TERRACOTTA = RGBColor(0xB8, 0x57, 0x41)
CREAM = RGBColor(0xED, 0xE3, 0xDA)
DARK = RGBColor(0x31, 0x39, 0x4D)
GRAY = RGBColor(0x62, 0x6B, 0x73)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

CONSTRUCT_DEFS = {
    "egocentric_encoding": (
        "Egocentric encoding",
        "Relation of a visible object to the viewer's current pose.",
        "Read the direction from the camera/viewer frame in the single current image.",
    ),
    "allocentric_encoding": (
        "Allocentric encoding",
        "Relation between two visible objects in a trusted object-centered frame.",
        "Requires intrinsic object facing (object-frame edges); viewer-relative approximations are rejected.",
    ),
    "spatial_working_memory": (
        "Spatial working memory",
        "Recall an encoded relation after the object leaves view and navigation continues.",
        "Encode the object in the first frame, follow the ordered navigation frames, then recall the earlier relation.",
    ),
    "invisible_displacement": (
        "Invisible displacement",
        "Track an object's location after it is hidden and relocated out of view.",
        "The object is visible before occlusion, moves while hidden, and remains absent in the final view.",
    ),
    "spatial_updating": (
        "Spatial updating",
        "Update the bearing of a static object after the viewer moves.",
        "Read the ordered navigation sequence: encode the object early, track your own motion through the intermediate frames, then report the updated bearing at the final pose.",
    ),
    "perspective_taking": (
        "Perspective-taking",
        "Imagine standing at landmark A facing B; report where landmark C is.",
        "Heading is relational (A→B) from positions — no object-intrinsic front required.",
    ),
    "route_knowledge": (
        "Route knowledge",
        "Choose the matching turn sequence for a walked source→goal route.",
        "MCQ over derive_turns() on the graph-snapped trajectory; not full-episode recall.",
    ),
    "survey_based_route_planning": (
        "Survey-based route planning",
        "Judge layout direction/distance (or first-hop under a recorded passage closure).",
        "Path must be untraversed; answers are not turn sequences (those are route_knowledge).",
    ),
}


def _set_run(run, text: str, *, size: int, bold: bool = False, color=DARK) -> None:
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    *,
    size=14,
    bold=False,
    color=DARK,
    align=PP_ALIGN.LEFT,
    lines: list[tuple[str, dict]] | None = None,
    text: str | None = None,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if lines is None:
        lines = [(text or "", {"size": size, "bold": bold, "color": color})]
    first = True
    for content, style in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = style.get("align", align)
        p.space_after = Pt(style.get("space_after", 4))
        run = p.add_run()
        _set_run(
            run,
            content,
            size=style.get("size", size),
            bold=style.get("bold", bold),
            color=style.get("color", color),
        )
    return box


def add_rect(slide, left, top, width, height, fill: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def _picture_size(path: str, max_w: int, max_h: int) -> tuple[int, int]:
    """Return EMU width/height fitting inside max box, preserving aspect."""
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    return int(iw * scale), int(ih * scale)


def clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)


def blank_layout(prs: Presentation):
    return prs.slide_layouts[10]  # BLANK


def load_items(paths: list[Path]) -> list[dict]:
    items = []
    for path in paths:
        if not path.exists():
            continue
        data = json.loads(path.read_text())
        for item in data.get("items", []):
            enriched = dict(item)
            enriched["_draft_source"] = path.stem
            items.append(enriched)
    if not items:
        joined = ", ".join(str(path) for path in paths)
        raise FileNotFoundError(f"No draft item JSON found: {joined}")
    return items


def select_examples(items: list[dict], construct: str, limit: int) -> list[dict]:
    """Choose distinct, valid concise examples in input-file order."""
    candidates = [
        item
        for item in items
        if item.get("construct") == construct
        and item.get("status") == "ok"
        and item.get("question")
        and item.get("question_style") == "concise"
        and all(Path(path).is_file() for path in (item.get("image_paths") or []))
    ]
    selected = []
    seen_pairs = set()
    for item in candidates:
        pair_key = item.get("paired_item_id") or item.get("item_id")
        if pair_key in seen_pairs:
            continue
        seen_pairs.add(pair_key)
        selected.append(item)
        if len(selected) >= limit:
            break
    return selected


def unsupported_reason(items: list[dict], construct: str) -> str:
    valid_without_images = [
        item
        for item in items
        if item.get("construct") == construct
        and item.get("status") == "ok"
        and item.get("question_style") == "concise"
    ]
    if valid_without_images:
        return (
            "raw image files referenced by the draft are unavailable; regenerate "
            "the draft from a current episode export"
        )
    for item in items:
        if item.get("construct") == construct and item.get("status") == "unsupported":
            rationale = item.get("distractor_rationale") or {}
            return str(rationale.get("reason") or "insufficient metadata")
    return "No strict, GT-provable example is available in the supplied drafts."


def refresh_progress_slide(slide) -> None:
    clear_slide(slide)
    W = Inches(10)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.7), NAVY)
    add_textbox(
        slide,
        Inches(0.4),
        Inches(0.15),
        Inches(9),
        Inches(0.45),
        text="Advances so far (generation pipeline)",
        size=22,
        bold=True,
        color=WHITE,
    )
    bullets = [
        "Taxonomy locked: 4 classes × 8 constructs (+ FoR axis)",
        "SPOC episode → Episode GT (poses, visibility, edges, displacement, layout)",
        "First-draft MC items: CODE-locked answers + answer_source from metadata",
        "Frame annotator for GT review (not model input)",
        "Still upcoming: GT Validator · vision-necessity · FREEZE · Model Evaluation",
    ]
    lines = [(f"•  {b}", {"size": 15, "bold": False, "color": DARK, "space_after": 10}) for b in bullets]
    add_textbox(slide, Inches(0.5), Inches(1.0), Inches(9), Inches(3.8), lines=lines)
    add_textbox(
        slide,
        Inches(0.5),
        Inches(5.05),
        Inches(9),
        Inches(0.35),
        text="Not a frozen eval set — draft candidates only.",
        size=12,
        bold=True,
        color=TERRACOTTA,
    )


def add_divider(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), NAVY)
    add_textbox(
        slide,
        Inches(0.6),
        Inches(1.8),
        Inches(8.8),
        Inches(0.8),
        text=title,
        size=32,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(2.8),
        Inches(8.4),
        Inches(1.0),
        text=subtitle,
        size=16,
        color=CREAM,
        align=PP_ALIGN.CENTER,
    )


def add_construct_header(
    prs: Presentation,
    name: str,
    definition: str,
    how_to_read: str,
    *,
    n_examples: int,
) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.7), TEAL)
    add_textbox(
        slide,
        Inches(0.4),
        Inches(0.15),
        Inches(9),
        Inches(0.45),
        text=name,
        size=22,
        bold=True,
        color=WHITE,
    )
    add_textbox(
        slide,
        Inches(0.5),
        Inches(0.8),
        Inches(9),
        Inches(1.2),
        text=definition,
        size=18,
        color=DARK,
    )
    add_textbox(
        slide,
        Inches(0.5),
        Inches(2.25),
        Inches(9),
        Inches(2.0),
        lines=[
            ("How to read the example", {"size": 14, "bold": True, "color": TEAL}),
            (how_to_read, {"size": 15, "color": DARK, "space_after": 12}),
            (
                f"{n_examples} strict GT-backed example(s) selected automatically from the supplied drafts."
                if n_examples
                else "No strict GT-backed example is available; the next slide explains the blocker.",
                {"size": 13, "color": GRAY},
            ),
            ("Answers are locked by code from episode GT — not invented by an LLM.", {"size": 14, "color": GRAY}),
        ],
    )


def _fit_images(slide, image_paths: list[str], left, top, max_w, max_h) -> None:
    n = len(image_paths)
    if n == 0:
        return
    gap = Inches(0.12)
    label_h = Inches(0.22)
    if n == 1:
        w, h = _picture_size(image_paths[0], int(max_w), int(max_h - label_h))
        slide.shapes.add_picture(image_paths[0], left, top, width=w, height=h)
        return
    each_w = Emu(int((max_w - gap) / 2))
    labels = ["encoding / earlier", "query / final"]
    for i, path in enumerate(image_paths[:2]):
        x = left + i * (each_w + gap)
        add_textbox(
            slide,
            x,
            top,
            each_w,
            label_h,
            text=labels[i],
            size=10,
            bold=True,
            color=TEAL,
            align=PP_ALIGN.CENTER,
        )
        w, h = _picture_size(path, int(each_w), int(max_h - label_h))
        slide.shapes.add_picture(path, x, top + label_h, width=w, height=h)


def _frame_labels(item: dict) -> list[str]:
    paths = item.get("image_paths") or []
    trajectory = item.get("agent_trajectory") or []
    labels = []
    for index, path in enumerate(paths):
        step = None
        if index < len(trajectory) and trajectory[index].get("image_path") == path:
            step = trajectory[index].get("step")
        role = ""
        if index == 0:
            role = " · encode"
        elif index == len(paths) - 1:
            role = " · query"
        labels.append(f"t={step if step is not None else index}{role}")
    return labels


def add_sequence_slides(
    prs: Presentation,
    item: dict,
    construct_title: str,
    example_n: int,
    *,
    frames_per_slide: int = 6,
) -> None:
    paths = list(item.get("image_paths") or [])
    if len(paths) <= 2:
        return
    labels = _frame_labels(item)
    total = (len(paths) + frames_per_slide - 1) // frames_per_slide
    for page, offset in enumerate(range(0, len(paths), frames_per_slide), start=1):
        slide = prs.slides.add_slide(blank_layout(prs))
        add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.55), NAVY)
        add_textbox(
            slide,
            Inches(0.3),
            Inches(0.1),
            Inches(9.4),
            Inches(0.4),
            text=(
                f"{construct_title} · Example {example_n} · ordered sequence "
                f"{page}/{total}"
            ),
            size=16,
            bold=True,
            color=WHITE,
        )
        add_textbox(
            slide,
            Inches(0.35),
            Inches(0.63),
            Inches(9.3),
            Inches(0.35),
            text=(
                "Read left→right, top→bottom. The first frame encodes the fact; "
                "intermediate frames create the navigation delay; the last frame is the query."
            ),
            size=11,
            color=GRAY,
            align=PP_ALIGN.CENTER,
        )

        chunk = paths[offset : offset + frames_per_slide]
        for local_index, path in enumerate(chunk):
            if not Path(path).exists():
                raise FileNotFoundError(path)
            row, col = divmod(local_index, 3)
            x = Inches(0.3 + col * 3.25)
            y = Inches(1.08 + row * 2.15)
            cell_w, cell_h = Inches(3.05), Inches(1.72)
            add_textbox(
                slide,
                x,
                y,
                cell_w,
                Inches(0.25),
                text=labels[offset + local_index],
                size=10,
                bold=offset + local_index in (0, len(paths) - 1),
                color=TEAL,
                align=PP_ALIGN.CENTER,
            )
            w, h = _picture_size(path, int(cell_w), int(cell_h - Inches(0.25)))
            slide.shapes.add_picture(
                path,
                x + Emu(int((cell_w - w) / 2)),
                y + Inches(0.27),
                width=w,
                height=h,
            )

        k = item.get("difficulty")
        add_textbox(
            slide,
            Inches(0.35),
            Inches(5.28),
            Inches(9.3),
            Inches(0.22),
            text=(
                f"Sequence: {len(paths)} frames"
                + (f" · delay k={k} navigation steps" if k is not None else "")
                + " · raw images shown to the model"
            ),
            size=9,
            color=GRAY,
            align=PP_ALIGN.CENTER,
        )


def add_example_slide(
    prs: Presentation, item: dict, construct_title: str, example_n: int
) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.55), NAVY)
    add_textbox(
        slide,
        Inches(0.3),
        Inches(0.1),
        Inches(9.4),
        Inches(0.4),
        text=f"{construct_title}  ·  Example {example_n} · question and answer",
        size=16,
        bold=True,
        color=WHITE,
    )

    paths = list(item.get("image_paths") or [])
    for p in paths:
        if not Path(p).exists():
            raise FileNotFoundError(p)

    # Long sequences receive one or more preceding timeline slides. The Q&A
    # slide repeats only the encoding and query endpoints as a compact recap.
    display_paths = paths if len(paths) <= 2 else [paths[0], paths[-1]]
    _fit_images(slide, display_paths, Inches(0.25), Inches(0.7), Inches(5.0), Inches(4.0))
    if len(paths) > 2:
        add_textbox(
            slide,
            Inches(0.35),
            Inches(4.72),
            Inches(4.8),
            Inches(0.35),
            text=f"Endpoint recap · full {len(paths)}-frame sequence on preceding slide(s)",
            size=9,
            color=GRAY,
            align=PP_ALIGN.CENTER,
        )

    # Q&A right
    rx, rw = Inches(5.4), Inches(4.3)
    add_textbox(slide, rx, Inches(0.7), rw, Inches(0.3), text="Question", size=11, bold=True, color=TEAL)
    q = item.get("question") or ""
    # shorten verbose for slide readability
    if len(q) > 420:
        q = q[:400].rsplit(" ", 1)[0] + "…"
    add_textbox(slide, rx, Inches(0.95), rw, Inches(1.7), text=q, size=11, color=DARK)

    opts = item.get("options") or {}
    ans = item.get("answer")
    opt_lines = []
    for key in sorted(opts.keys()):
        label = f"{key}.  {opts[key]}"
        if key == ans:
            opt_lines.append((label, {"size": 12, "bold": True, "color": TEAL, "space_after": 6}))
        else:
            opt_lines.append((label, {"size": 12, "bold": False, "color": DARK, "space_after": 6}))
    add_textbox(slide, rx, Inches(2.7), rw, Inches(0.25), text="Options", size=11, bold=True, color=TEAL)
    add_textbox(slide, rx, Inches(2.95), rw, Inches(1.6), lines=opt_lines)

    ans_text = opts.get(ans, "")
    add_textbox(
        slide,
        rx,
        Inches(4.55),
        rw,
        Inches(0.4),
        text=f"Answer: {ans} — {ans_text}",
        size=13,
        bold=True,
        color=TEAL,
    )

    for_ = item.get("frame_of_reference", "")
    style = item.get("question_style", "")
    iid = item.get("item_id", "")
    short_id = iid[:60]
    scene = item.get("scene_id", "")
    encoding = item.get("encoding_step")
    query = item.get("query_step")
    add_textbox(
        slide,
        Inches(0.25),
        Inches(5.2),
        Inches(9.5),
        Inches(0.3),
        text=(
            f"{scene} · FoR={for_} · steps={encoding}→{query} · "
            f"style={style} · {short_id}"
        ),
        size=9,
        color=GRAY,
    )


def add_unsupported_slide(
    prs: Presentation,
    construct_title: str,
    definition: str,
    reason: str,
) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.7), TERRACOTTA)
    add_textbox(
        slide,
        Inches(0.4),
        Inches(0.15),
        Inches(9),
        Inches(0.45),
        text=f"{construct_title} · strict example unavailable",
        size=22,
        bold=True,
        color=WHITE,
    )
    add_textbox(
        slide,
        Inches(0.5),
        Inches(1.1),
        Inches(9),
        Inches(3.5),
        lines=[
            (
                f"Definition: {definition}",
                {"size": 16, "color": DARK, "space_after": 14},
            ),
            (
                f"Metadata blocker: {reason.replace('_', ' ')}.",
                {"size": 15, "bold": True, "color": TERRACOTTA, "space_after": 14},
            ),
            (
                "No example is shown because an unprovable draft would violate the benchmark's GT-only invariant.",
                {"size": 15, "color": DARK, "space_after": 10},
            ),
            (
                "This is an explicit unsupported result, not a missing slide or generation failure.",
                {"size": 14, "color": GRAY},
            ),
        ],
    )


def add_final_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.7), NAVY)
    add_textbox(
        slide,
        Inches(0.4),
        Inches(0.15),
        Inches(9),
        Inches(0.45),
        text="Takeaways",
        size=22,
        bold=True,
        color=WHITE,
    )
    add_textbox(
        slide,
        Inches(0.5),
        Inches(1.0),
        Inches(9),
        Inches(4.0),
        lines=[
            ("• Taxonomy is the capability axis; FoR is cross-cutting.", {"size": 15, "space_after": 10}),
            ("• Each construct is shown with strict GT-backed examples or an explicit metadata blocker.", {"size": 15, "space_after": 10}),
            ("• Answers are CODE-locked to GT; images shown = model input.", {"size": 15, "space_after": 10}),
            ("• Temporal examples show the ordered frame sequence, not only its endpoints.", {"size": 15, "space_after": 10}),
            ("• Unsupported constructs remain visible as honest metadata gaps.", {"size": 15, "space_after": 10}),
            ("• Next wall: validate → vision-necessity → FREEZE → evaluate VLMs.", {"size": 15, "bold": True, "color": TEAL}),
        ],
    )


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Build understandable, per-construct benchmark example slides."
    )
    parser.add_argument("--template", type=Path, default=TEMPLATE)
    parser.add_argument("--output", type=Path, default=OUTPUT)
    parser.add_argument(
        "--draft-json",
        type=Path,
        action="append",
        dest="draft_jsons",
        help="Draft JSON input; repeat to combine episodes (defaults to known local drafts)",
    )
    parser.add_argument("--examples-per-construct", type=int, default=2)
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    draft_jsons = args.draft_jsons or DEFAULT_DRAFT_JSONS
    if not args.template.exists():
        raise SystemExit(f"Template missing: {args.template}")
    if args.examples_per_construct < 1:
        raise SystemExit("--examples-per-construct must be at least 1")
    args.output.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(args.template, args.output)

    prs = Presentation(str(args.output))
    items = load_items(draft_jsons)

    # Slide 8 (index 7): refresh progress
    refresh_progress_slide(prs.slides[7])

    add_divider(
        prs,
        "First-draft items — what the model sees",
        "Strict GT-backed MC candidates · raw ordered frames · not a frozen set",
    )

    for construct, (title, definition, how_to_read) in CONSTRUCT_DEFS.items():
        examples = select_examples(items, construct, args.examples_per_construct)
        add_construct_header(
            prs,
            title,
            definition,
            how_to_read,
            n_examples=len(examples),
        )
        if not examples:
            add_unsupported_slide(
                prs, title, definition, unsupported_reason(items, construct)
            )
            continue
        for n, item in enumerate(examples, start=1):
            add_sequence_slides(prs, item, title, n)
            add_example_slide(prs, item, title, n)

    add_final_slide(prs)
    prs.save(str(args.output))
    print(f"Wrote {args.output} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
